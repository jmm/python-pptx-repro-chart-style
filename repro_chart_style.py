import pptx
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

presentation = pptx.Presentation("repro-chart-style-template.pptx")

chart_data = ChartData()
chart_data.categories = ['Yes', 'No']
chart_data.add_series('Series 1', (42, 24))

for chart_style in range(1, 49):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    placeholder = slide.placeholders[13]
    chart = placeholder.insert_chart(XL_CHART_TYPE.DOUGHNUT, chart_data).chart
    chart.chart_style = chart_style

presentation.save("repro-chart-style.pptx")
